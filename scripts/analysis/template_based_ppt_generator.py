#!/usr/bin/env python3
"""
Template-Based PowerPoint Generator
Uses existing Synoptek template formatting and style
"""

import json
import os
import pandas as pd
import numpy as np
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import shutil

class TemplateBasedPowerPointGenerator:
    def __init__(self):
        self.ai_data_file = "reports/current/ai_enhanced_industry_analysis_20250725.json"
        self.output_dir = "reports/current/presentation_style_analysis_20250725"
        self.template_file = "templates/SynoptekTransitionPlanJuly2025.pptx"
        self.ppt_file = f"{self.output_dir}/synoptek_charge_assessment_template_based.pptx"
        
    def load_ai_data(self):
        """Load the AI-enhanced analysis data."""
        if not os.path.exists(self.ai_data_file):
            print(f"Error: AI-enhanced data file not found: {self.ai_data_file}")
            return None

        with open(self.ai_data_file, 'r') as f:
            return json.load(f)
    
    def create_presentation_analysis(self, data):
        """Create presentation-style analysis focusing on charge assessment."""
        benchmarks = data.get('benchmarks', [])
        
        # Filter for Synoptek records
        synoptek_records = [b for b in benchmarks if b.get('vendor', '').lower() == 'synoptek']
        
        # Create DataFrame
        df = pd.DataFrame(synoptek_records)
        
        if df.empty:
            print("No Synoptek records found in the dataset.")
            return None
        
        # Focus on charge assessment
        df['actual_cost'] = df['actual_spend']
        df['subcategory'] = df['subcategory'].fillna('Unknown')
        df['category'] = df['category'].fillna('Unknown')
        
        # Extract service information
        df['service_type'] = df['ai_categorization'].apply(
            lambda x: x.get('service_type', 'Unknown') if isinstance(x, dict) else 'Unknown'
        )
        
        df['primary_category'] = df['ai_categorization'].apply(
            lambda x: x.get('primary_category', 'Unknown') if isinstance(x, dict) else 'Unknown'
        )
        
        # Calculate charge assessment metrics
        df['service_identifier'] = df['subcategory'] + '_' + df['service_type']
        
        # Analyze charges by service type
        service_analysis = df.groupby('service_identifier').agg({
            'actual_cost': ['sum', 'mean', 'count', 'std'],
            'subcategory': 'first',
            'service_type': 'first',
            'primary_category': 'first'
        }).round(2)
        
        service_analysis.columns = ['Total_Cost', 'Avg_Monthly_Cost', 'Billing_Months', 'Cost_Std_Dev', 
                                 'Subcategory', 'Service_Type', 'Primary_Category']
        
        # Calculate charge assessment indicators
        service_analysis['Cost_Per_Service'] = service_analysis['Avg_Monthly_Cost']
        service_analysis['Cost_Variance'] = service_analysis['Cost_Std_Dev'] / service_analysis['Avg_Monthly_Cost']
        
        # Identify potential overcharges
        cost_threshold = service_analysis['Avg_Monthly_Cost'].quantile(0.9)  # Top 10%
        service_analysis['High_Cost_Flag'] = service_analysis['Avg_Monthly_Cost'] > cost_threshold
        
        # Calculate industry comparison metrics
        service_analysis['Industry_Comparison'] = service_analysis['Avg_Monthly_Cost'].apply(
            lambda x: 'Above Average' if x > cost_threshold else 'Normal' if x > service_analysis['Avg_Monthly_Cost'].median() else 'Below Average'
        )
        
        return df, service_analysis
    
    def assess_charge_fairness(self, df, service_analysis):
        """Assess whether charges are fair or overpriced."""
        
        # Calculate overall metrics
        total_cost = df['actual_cost'].sum()
        avg_cost_per_service = service_analysis['Avg_Monthly_Cost'].mean()
        median_cost = service_analysis['Avg_Monthly_Cost'].median()
        
        # Identify overpriced services
        overpriced_services = service_analysis[service_analysis['High_Cost_Flag']].copy()
        overpriced_cost = overpriced_services['Total_Cost'].sum()
        overpriced_percentage = (overpriced_cost / total_cost) * 100
        
        # Calculate fairness metrics
        cost_distribution = {
            'Normal': len(service_analysis[service_analysis['Industry_Comparison'] == 'Normal']),
            'Above_Average': len(service_analysis[service_analysis['Industry_Comparison'] == 'Above Average']),
            'Below_Average': len(service_analysis[service_analysis['Industry_Comparison'] == 'Below Average'])
        }
        
        # Assess overall fairness
        if overpriced_percentage > 30:
            overall_assessment = "POTENTIALLY OVERCHARGED"
            assessment_color = "red"
        elif overpriced_percentage > 15:
            overall_assessment = "MIXED - SOME CONCERNS"
            assessment_color = "orange"
        else:
            overall_assessment = "CHARGES APPEAR NORMAL"
            assessment_color = "green"
        
        assessment = {
            'total_cost': total_cost,
            'avg_cost_per_service': avg_cost_per_service,
            'median_cost': median_cost,
            'overpriced_services_count': len(overpriced_services),
            'overpriced_cost': overpriced_cost,
            'overpriced_percentage': overpriced_percentage,
            'cost_distribution': cost_distribution,
            'overall_assessment': overall_assessment,
            'assessment_color': assessment_color,
            'overpriced_services': overpriced_services,
            'service_analysis': service_analysis
        }
        
        return assessment
    
    def create_template_based_presentation(self, assessment):
        """Create presentation using the existing template."""
        
        # Check if template exists
        if not os.path.exists(self.template_file):
            print(f"Template file not found: {self.template_file}")
            print("Creating basic presentation instead...")
            return self.create_basic_presentation(assessment)
        
        try:
            # Copy template to output location
            shutil.copy2(self.template_file, self.ppt_file)
            
            # Load the copied template
            prs = Presentation(self.ppt_file)
            
            # Clear existing slides (keep first slide as template)
            while len(prs.slides) > 1:
                rId = prs.slides._sldIdLst[1].rId
                prs.part.drop_rel(rId)
                prs.slides._sldIdLst.remove(prs.slides._sldIdLst[1])
            
            # Update the title slide (first slide)
            if len(prs.slides) > 0:
                slide = prs.slides[0]
                if slide.shapes.title:
                    slide.shapes.title.text = "Synoptek Charge Assessment"
                # Update subtitle if it exists
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                        if "Transition" in shape.text or "Plan" in shape.text:
                            shape.text = f"Are We Being Overcharged?\n\nAssessment: {assessment['overall_assessment']}\nGenerated: {datetime.now().strftime('%Y-%m-%d')}"
            
            # Add new slides using template layouts
            self.add_executive_summary_slide(prs, assessment)
            self.add_charge_breakdown_slide(prs, assessment)
            self.add_overpriced_services_slide(prs, assessment)
            self.add_service_analysis_slide(prs, assessment)
            self.add_findings_slide(prs, assessment)
            self.add_strategic_recommendations_slide(prs, assessment)
            self.add_conclusion_slide(prs, assessment)
            
            # Save presentation
            prs.save(self.ppt_file)
            
            return True
            
        except Exception as e:
            print(f"Error using template: {e}")
            print("Falling back to basic presentation...")
            return self.create_basic_presentation(assessment)
    
    def create_basic_presentation(self, assessment):
        """Create basic presentation if template not available."""
        prs = Presentation()
        
        # Set slide size to 16:9
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Create slides
        self.add_title_slide(prs, assessment)
        self.add_executive_summary_slide(prs, assessment)
        self.add_charge_breakdown_slide(prs, assessment)
        self.add_overpriced_services_slide(prs, assessment)
        self.add_service_analysis_slide(prs, assessment)
        self.add_findings_slide(prs, assessment)
        self.add_strategic_recommendations_slide(prs, assessment)
        self.add_conclusion_slide(prs, assessment)
        
        # Save presentation
        prs.save(self.ppt_file)
        
        return True
    
    def add_title_slide(self, prs, assessment):
        """Add title slide."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        if title:
            title.text = "Synoptek Charge Assessment"
            title.text_frame.paragraphs[0].font.size = Pt(44)
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        
        # Set subtitle
        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = f"Are We Being Overcharged?\n\nAssessment: {assessment['overall_assessment']}\nGenerated: {datetime.now().strftime('%Y-%m-%d')}"
            subtitle.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
        
        return slide
    
    def add_executive_summary_slide(self, prs, assessment):
        """Add executive summary slide."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        if title:
            title.text = "Executive Summary"
            title.text_frame.paragraphs[0].font.size = Pt(36)
            title.text_frame.paragraphs[0].font.bold = True
        
        # Add content
        if len(slide.placeholders) > 1:
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.clear()
            
            # Add key metrics
            p = tf.add_paragraph()
            p.text = f"OVERALL ASSESSMENT: {assessment['overall_assessment']}"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(220, 20, 60) if assessment['overall_assessment'] == "POTENTIALLY OVERCHARGED" else RGBColor(255, 140, 0)
            
            tf.add_paragraph()
            
            # Add key metrics
            metrics = [
                f"Total Cost: ${assessment['total_cost']:,.2f}",
                f"Overpriced Services: {assessment['overpriced_services_count']} services flagged",
                f"Overpriced Amount: ${assessment['overpriced_cost']:,.2f}",
                f"Overpriced Percentage: {assessment['overpriced_percentage']:.1f}% of total spend"
            ]
            
            for metric in metrics:
                p = tf.add_paragraph()
                p.text = f"• {metric}"
                p.font.size = Pt(16)
                p.level = 1
        
        return slide
    
    def add_charge_breakdown_slide(self, prs, assessment):
        """Add charge breakdown slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        if title:
            title.text = "Charge Assessment Breakdown"
            title.text_frame.paragraphs[0].font.size = Pt(36)
            title.text_frame.paragraphs[0].font.bold = True
        
        # Add content
        if len(slide.placeholders) > 1:
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.clear()
            
            # Add cost distribution
            p = tf.add_paragraph()
            p.text = "Cost Distribution Analysis:"
            p.font.size = Pt(18)
            p.font.bold = True
            
            distribution = assessment['cost_distribution']
            tf.add_paragraph()
            
            for category, count in distribution.items():
                p = tf.add_paragraph()
                p.text = f"• {category.replace('_', ' ')}: {count} services"
                p.font.size = Pt(14)
                p.level = 1
            
            tf.add_paragraph()
            
            # Add industry comparison
            p = tf.add_paragraph()
            p.text = "Industry Comparison:"
            p.font.size = Pt(18)
            p.font.bold = True
            
            tf.add_paragraph()
            
            p = tf.add_paragraph()
            p.text = f"• Average Cost per Service: ${assessment['avg_cost_per_service']:,.2f}"
            p.font.size = Pt(14)
            p.level = 1
            
            p = tf.add_paragraph()
            p.text = f"• Median Cost per Service: ${assessment['median_cost']:,.2f}"
            p.font.size = Pt(14)
            p.level = 1
        
        return slide
    
    def add_overpriced_services_slide(self, prs, assessment):
        """Add overpriced services slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        if title:
            title.text = "Potentially Overpriced Services"
            title.text_frame.paragraphs[0].font.size = Pt(36)
            title.text_frame.paragraphs[0].font.bold = True
        
        # Add content
        if len(slide.placeholders) > 1:
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.clear()
            
            # Add top overpriced services
            top_overpriced = assessment['overpriced_services'].nlargest(8, 'Avg_Monthly_Cost')
            
            for idx, row in top_overpriced.iterrows():
                p = tf.add_paragraph()
                p.text = f"• {row['Subcategory']} ({row['Service_Type']})"
                p.font.size = Pt(14)
                p.font.bold = True
                
                p = tf.add_paragraph()
                p.text = f"  Monthly Cost: ${row['Avg_Monthly_Cost']:,.2f}"
                p.font.size = Pt(12)
                p.level = 1
                
                p = tf.add_paragraph()
                p.text = f"  Total Cost: ${row['Total_Cost']:,.2f}"
                p.font.size = Pt(12)
                p.level = 1
        
        return slide
    
    def add_service_analysis_slide(self, prs, assessment):
        """Add service type analysis slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        if title:
            title.text = "Service Type Analysis"
            title.text_frame.paragraphs[0].font.size = Pt(36)
            title.text_frame.paragraphs[0].font.bold = True
        
        # Add content
        if len(slide.placeholders) > 1:
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.clear()
            
            # Add highest cost service types
            service_costs = assessment['service_analysis'].groupby('Service_Type')['Total_Cost'].sum().sort_values(ascending=False).head(8)
            
            for service_type, cost in service_costs.items():
                service_count = len(assessment['service_analysis'][assessment['service_analysis']['Service_Type'] == service_type])
                avg_cost = assessment['service_analysis'][assessment['service_analysis']['Service_Type'] == service_type]['Avg_Monthly_Cost'].mean()
                
                p = tf.add_paragraph()
                p.text = f"• {service_type}"
                p.font.size = Pt(14)
                p.font.bold = True
                
                p = tf.add_paragraph()
                p.text = f"  Total Cost: ${cost:,.2f}"
                p.font.size = Pt(12)
                p.level = 1
                
                p = tf.add_paragraph()
                p.text = f"  Service Count: {service_count}"
                p.font.size = Pt(12)
                p.level = 1
                
                p = tf.add_paragraph()
                p.text = f"  Average Monthly Cost: ${avg_cost:,.2f}"
                p.font.size = Pt(12)
                p.level = 1
        
        return slide
    
    def add_findings_slide(self, prs, assessment):
        """Add findings and recommendations slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        if title:
            title.text = "Key Findings & Recommendations"
            title.text_frame.paragraphs[0].font.size = Pt(36)
            title.text_frame.paragraphs[0].font.bold = True
        
        # Add content
        if len(slide.placeholders) > 1:
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.clear()
            
            # Add findings based on assessment
            if assessment['overpriced_percentage'] > 30:
                p = tf.add_paragraph()
                p.text = "🚨 CRITICAL FINDINGS:"
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = RGBColor(220, 20, 60)
                
                findings = [
                    f"{assessment['overpriced_percentage']:.1f}% of total spend appears overpriced",
                    f"${assessment['overpriced_cost']:,.2f} in potentially overcharged costs",
                    f"{assessment['overpriced_services_count']} services require immediate review"
                ]
            elif assessment['overpriced_percentage'] > 15:
                p = tf.add_paragraph()
                p.text = "⚠️ MIXED FINDINGS:"
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 140, 0)
                
                findings = [
                    f"{assessment['overpriced_percentage']:.1f}% of total spend shows concerns",
                    f"${assessment['overpriced_cost']:,.2f} in potentially overcharged costs",
                    f"{assessment['overpriced_services_count']} services need review"
                ]
            else:
                p = tf.add_paragraph()
                p.text = "✅ POSITIVE FINDINGS:"
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 128, 0)
                
                findings = [
                    f"Only {assessment['overpriced_percentage']:.1f}% of total spend shows concerns",
                    f"${assessment['overpriced_cost']:,.2f} in potentially overcharged costs",
                    f"{assessment['overpriced_services_count']} services flagged for review"
                ]
            
            for finding in findings:
                p = tf.add_paragraph()
                p.text = f"• {finding}"
                p.font.size = Pt(14)
                p.level = 1
            
            tf.add_paragraph()
            
            # Add recommendations
            p = tf.add_paragraph()
            p.text = "RECOMMENDATIONS:"
            p.font.size = Pt(18)
            p.font.bold = True
            
            recommendations = [
                "Immediate Action: Review all flagged services",
                "Negotiation: Use analysis to negotiate with Synoptek",
                "Alternative Quotes: Obtain competitive quotes",
                "Contract Review: Reassess current contract terms"
            ]
            
            for rec in recommendations:
                p = tf.add_paragraph()
                p.text = f"• {rec}"
                p.font.size = Pt(14)
                p.level = 1
        
        return slide
    
    def add_strategic_recommendations_slide(self, prs, assessment):
        """Add strategic recommendations slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        if title:
            title.text = "Strategic Recommendations"
            title.text_frame.paragraphs[0].font.size = Pt(36)
            title.text_frame.paragraphs[0].font.bold = True
        
        # Add content
        if len(slide.placeholders) > 1:
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.clear()
            
            # Immediate actions
            p = tf.add_paragraph()
            p.text = "Immediate Actions (Next 30 Days):"
            p.font.size = Pt(16)
            p.font.bold = True
            
            immediate_actions = [
                f"Service Review: Audit all {assessment['overpriced_services_count']} flagged services",
                "Market Comparison: Compare pricing with industry benchmarks",
                "Vendor Discussion: Schedule meeting with Synoptek",
                "Alternative Quotes: Obtain competitive quotes"
            ]
            
            for action in immediate_actions:
                p = tf.add_paragraph()
                p.text = f"• {action}"
                p.font.size = Pt(12)
                p.level = 1
            
            tf.add_paragraph()
            
            # Short-term actions
            p = tf.add_paragraph()
            p.text = "Short-term Actions (Next 90 Days):"
            p.font.size = Pt(16)
            p.font.bold = True
            
            short_term_actions = [
                "Contract Negotiation: Use analysis to negotiate better rates",
                "Service Optimization: Consolidate or eliminate overpriced services",
                "Cost Monitoring: Implement regular cost review process",
                "Performance Tracking: Monitor service quality vs cost"
            ]
            
            for action in short_term_actions:
                p = tf.add_paragraph()
                p.text = f"• {action}"
                p.font.size = Pt(12)
                p.level = 1
        
        return slide
    
    def add_conclusion_slide(self, prs, assessment):
        """Add conclusion slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        if title:
            title.text = "Conclusion"
            title.text_frame.paragraphs[0].font.size = Pt(36)
            title.text_frame.paragraphs[0].font.bold = True
        
        # Add content
        if len(slide.placeholders) > 1:
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.clear()
            
            # Add conclusion
            p = tf.add_paragraph()
            p.text = f"Assessment: {assessment['overall_assessment']}"
            p.font.size = Pt(18)
            p.font.bold = True
            
            tf.add_paragraph()
            
            p = tf.add_paragraph()
            p.text = f"Primary Finding: {assessment['overpriced_percentage']:.1f}% of total spend ({assessment['overpriced_services_count']} services) appears potentially overpriced."
            p.font.size = Pt(16)
            
            tf.add_paragraph()
            
            if assessment['overpriced_percentage'] > 30:
                recommendation = "Immediate action required."
            elif assessment['overpriced_percentage'] > 15:
                recommendation = "Selective review recommended."
            else:
                recommendation = "Minor optimization opportunities exist."
            
            p = tf.add_paragraph()
            p.text = f"Recommendation: {recommendation}"
            p.font.size = Pt(16)
            p.font.bold = True
            
            tf.add_paragraph()
            
            # Add next steps
            p = tf.add_paragraph()
            p.text = "Next Steps:"
            p.font.size = Pt(16)
            p.font.bold = True
            
            next_steps = [
                "Review detailed CSV reports for granular analysis",
                "Schedule meeting with Synoptek to discuss findings",
                "Obtain competitive quotes for overpriced services",
                "Implement regular cost monitoring process"
            ]
            
            for step in next_steps:
                p = tf.add_paragraph()
                p.text = f"• {step}"
                p.font.size = Pt(14)
                p.level = 1
        
        return slide
    
    def generate_template_based_analysis(self):
        """Generate the complete template-based PowerPoint presentation."""
        print("=" * 70)
        print("    TEMPLATE-BASED POWERPOINT GENERATOR")
        print("=" * 70)
        print()
        
        # Load data
        data = self.load_ai_data()
        if not data:
            return False
        
        print("📊 Creating template-based PowerPoint presentation...")
        
        # Create analysis
        df, service_analysis = self.create_presentation_analysis(data)
        if df is None:
            print("❌ No Synoptek records found for analysis!")
            return False
        
        print(f"📋 Found {len(df)} Synoptek service records for assessment")
        
        # Assess charge fairness
        assessment = self.assess_charge_fairness(df, service_analysis)
        
        # Create PowerPoint presentation using template
        print("📈 Creating PowerPoint slides using template...")
        success = self.create_template_based_presentation(assessment)
        
        if success:
            print("✅ Template-based PowerPoint presentation generated successfully!")
            print(f"📁 Output file: {self.ppt_file}")
            print(f"📊 Presentation includes:")
            print(f"   - Title slide (using template formatting)")
            print(f"   - Executive summary")
            print(f"   - Charge breakdown")
            print(f"   - Overpriced services")
            print(f"   - Service analysis")
            print(f"   - Findings & recommendations")
            print(f"   - Strategic recommendations")
            print(f"   - Conclusion")
            
            print()
            print(f"🎯 ASSESSMENT: {assessment['overall_assessment']}")
            print(f"💰 Overpriced Amount: ${assessment['overpriced_cost']:,.2f}")
            print(f"📊 Overpriced Percentage: {assessment['overpriced_percentage']:.1f}%")
            
            return True
        else:
            print("❌ Template-based PowerPoint generation failed!")
            return False

def main():
    """Main function to generate template-based PowerPoint presentation."""
    generator = TemplateBasedPowerPointGenerator()
    success = generator.generate_template_based_analysis()
    
    if success:
        print()
        print("🎉 Template-based PowerPoint presentation completed successfully!")
        print("📋 Ready for executive presentation")
        print("💼 Use presentation for charge assessment and negotiations")
    else:
        print("❌ Template-based PowerPoint generation failed!")
        return 1
    
    return 0

if __name__ == "__main__":
    exit(main())